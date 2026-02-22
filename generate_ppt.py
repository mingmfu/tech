from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]  # 标题页
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景为深蓝色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)  # slate-900
    
    # 标题
    title_box = slide.shapes.title
    title_box.text = title
    title_para = title_box.text_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    
    # 副标题
    subtitle_box = slide.placeholders[1]
    subtitle_box.text = subtitle
    subtitle_para = subtitle_box.text_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(148, 163, 184)  # slate-400
    
    return slide

def add_content_slide(prs, title, content_lines):
    slide_layout = prs.slide_layouts[1]  # 标题和内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)
    
    # 标题
    title_box = slide.shapes.title
    title_box.text = title
    title_para = title_box.text_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    
    # 内容
    body_box = slide.placeholders[1]
    tf = body_box.text_frame
    tf.clear()
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(203, 213, 225)  # slate-300
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_content, right_title, right_content):
    slide_layout = prs.slide_layouts[5]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)
    
    # 标题
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    
    # 左侧栏
    left_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5))
    left_frame = left_shape.text_frame
    left_frame.word_wrap = True
    
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(59, 130, 246)  # blue-500
    
    for line in left_content:
        p = left_frame.add_paragraph()
        p.text = "• " + line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(203, 213, 225)
        p.space_after = Pt(8)
    
    # 右侧栏
    right_shape = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5))
    right_frame = right_shape.text_frame
    right_frame.word_wrap = True
    
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(139, 92, 246)  # purple-500
    
    for line in right_content:
        p = right_frame.add_paragraph()
        p.text = "• " + line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(203, 213, 225)
        p.space_after = Pt(8)
    
    return slide

# ===== 开始创建幻灯片 =====

# 第1页：封面
add_title_slide(prs, 
    "皮升级液体分配技术",
    "微观世界的精准操控革命\n2025年2月")

# 第2页：执行摘要
add_content_slide(prs, "执行摘要", [
    "皮升级（picoliter，10⁻¹²升）液体分配技术正成为精密制造和生命科学领域的核心使能技术",
    "",
    "核心突破：",
    "• 压电陶瓷驱动、微机电系统（MEMS）和智能反馈控制的融合创新",
    "• 实现±1%以内的体积精度和每秒数千次的分配频率",
    "",
    "市场规模：",
    "• 预计2028年将突破120亿美元，年复合增长率18%",
    "",
    "应用前景：",
    "• 正在重塑基因测序、类器官芯片、微流控诊断和先进封装等前沿领域"
])

# 第3页：什么是皮升级
add_content_slide(prs, "什么是皮升级？", [
    "定义：",
    "• 1 皮升（picoliter）= 10⁻¹² 升 = 0.000000000001 升",
    "• 相当于人类头发直径的 1/5",
    "• 约等于一个红细胞的体积",
    "",
    "技术挑战：",
    "• 传统液体分配依赖气动或活塞驱动，最小体积受限在微升级（10⁻⁶升）",
    "• 皮升级分配需要纳米级精度控制和毫秒级响应速度",
    "• 液滴直径仅为20-100μm，对材料和工艺提出极高要求",
    "",
    "技术意义：",
    "• 实现单细胞级别、单分子级别的精确操控",
    "• 支撑高通量筛选、微量反应体系、精准医疗等前沿应用"
])

# 第4页：技术突破1 - 压电陶瓷驱动
add_content_slide(prs, "技术突破一：压电陶瓷驱动", [
    "核心技术原理：",
    "• 利用压电材料的逆压电效应，通过精确控制电压实现纳米级位移",
    "• 配合快速响应阀门和高刚度喷嘴结构，毫秒级完成液滴生成和断裂",
    "",
    "关键性能指标：",
    "• 分配体积范围：10pL - 1000pL（可调）",
    "• 体积精度：CV < 1%（部分高端系统达0.5%）",
    "• 分配频率：最高10kHz",
    "• 液滴直径：20-100μm",
    "",
    "技术优势：",
    "• 响应速度快、控制精度高、重复性好",
    "• 可处理高粘度液体和含颗粒悬浮液",
    "• 无机械磨损，使用寿命长"
])

# 第5页：技术突破2 - MEMS微纳加工
add_content_slide(prs, "技术突破二：MEMS微纳加工", [
    "技术概述：",
    "• 微机电系统（MEMS）技术实现皮升级分配器的微型化和集成化",
    "• 采用硅基或玻璃基微加工，喷嘴直径可小至20μm",
    "",
    "喷嘴微结构优化：",
    "• 表面疏水/亲水图案化处理，精确控制液滴形态",
    "• 集成微加热器，实现高粘度液体的温控分配",
    "• 低表面能涂层防止液体粘附和交叉污染",
    "",
    "多通道并行设计：",
    "• 单个芯片集成数百至数千个独立分配通道",
    "• 通道间距可小至100μm，实现超高密度点阵打印",
    "• 独立寻址控制，每个通道可分配不同体积或不同液体"
])

# 第6页：技术突破3 - 智能反馈控制
add_content_slide(prs, "技术突破三：智能反馈控制", [
    "实时液滴监测技术：",
    "• 频闪成像系统：以微秒级曝光时间捕获液滴形成过程",
    "• 激光衍射传感器：非接触式实时测量液滴体积",
    "• 高速相机视觉检测：记录液滴形态和飞行轨迹",
    "",
    "AI反馈算法：",
    "• 基于深度学习的液滴形态预测和动态调整",
    "• 自动补偿液体粘度变化、温度漂移和喷嘴磨损",
    "• 闭环控制确保长期稳定的高精度分配",
    "",
    "智能化优势：",
    "• 预测喷嘴堵塞，实现自适应维护",
    "• 自动优化分配参数，零故障运行",
    "• 实时质量监控和数据追溯"
])

# 第7页：技术突破4 - 材料科学
add_content_slide(prs, "技术突破四：材料科学", [
    "低表面能涂层技术：",
    "• 喷嘴内壁采用氟碳聚合物或类金刚石碳（DLC）涂层",
    "• 降低液体粘附，防止堵塞和交叉污染",
    "• 提高液滴形成的一致性和稳定性",
    "",
    "陶瓷与复合材料喷嘴：",
    "• 氧化锆陶瓷喷嘴替代传统不锈钢",
    "• 更高的硬度、耐磨性和化学惰性",
    "• 使用寿命延长10倍以上",
    "",
    "生物相容性材料：",
    "• 开发全生物降解的分配系统",
    "• 无细胞毒性材料用于体内药物递送",
    "• 满足GMP规范的无菌和热原要求"
])

# 第8页：应用场景1 - 基因测序与合成生物学
add_content_slide(prs, "应用场景一：基因测序与合成生物学", [
    "高通量基因测序（NGS）样本制备：",
    "• 在微小反应室（pL级体积）中进行文库制备，试剂消耗降低1000倍",
    "• 单芯片可同时处理数万个样本，实现超高通量筛选",
    "• 用于单细胞测序，精确分配裂解缓冲液和反应酶",
    "",
    "DNA合成与编辑：",
    "• 电化学阵列合成：通过皮升级分配在电极表面逐层合成寡核苷酸",
    "• CRISPR文库构建：高通量分配guide RNA和Cas蛋白",
    "• 合成生物学：自动化构建基因回路和代谢通路",
    "",
    "市场规模：基因测序液体处理系统市场2027年预计达45亿美元，年复合增长率18%"
])

# 第9页：应用场景2 - 类器官与器官芯片
add_content_slide(prs, "应用场景二：类器官与器官芯片", [
    "类器官构建：",
    "• 皮升级分配实现单细胞或细胞小团的精确放置",
    "• 构建复杂多细胞结构的3D生物打印",
    "• 模拟肿瘤微环境，用于药物筛选和个性化医疗",
    "",
    "器官芯片（Organ-on-a-Chip）：",
    "• 微流控芯片中精确分配培养基、药物和细胞外基质",
    "• 模拟血管-组织界面的物质交换",
    "• 皮升级灌流系统实现动态培养环境的精确控制",
    "",
    "产业影响：",
    "• 器官芯片市场预计2030年达12亿美元",
    "• 皮升级分配是核心使能技术之一",
    "• 推动精准医疗和药物研发范式变革"
])

# 第10页：应用场景3 - 微流控诊断与POCT
add_content_slide(prs, "应用场景三：微流控诊断与POCT", [
    "数字PCR（dPCR）和单分子检测：",
    "• 将反应体系分割为数十万至数百万个皮升级微液滴",
    "• 每个液滴作为独立的PCR反应室",
    "• 实现绝对定量和极低丰度靶标的检测",
    "",
    "侧向层析和免疫检测：",
    "• 皮升级分配实现检测抗体的精确定量包被",
    "• 多重检测芯片上不同抗体的并行分配",
    "• 用于传染病快速检测、肿瘤标志物筛查",
    "",
    "POCT市场：",
    "• 受COVID-19推动，即时检测市场快速扩张",
    "• 皮升级液体处理系统需求激增",
    "• 从实验室走向家庭和社区检测"
])

# 第11页：应用场景4 - 先进封装与微电子
add_content_slide(prs, "应用场景四：先进封装与微电子", [
    "倒装芯片（Flip-Chip）底部填充：",
    "• 皮升级精度控制底部填充胶的分配量",
    "• 避免溢胶和填充不足，提高封装良率",
    "• 适用于高密度互连的微凸点封装",
    "",
    "微LED巨量转移：",
    "• 将微LED芯片从生长基板转移至显示背板",
    "• 皮升级粘胶的精确分配确保转移精度",
    "• 用于AR/VR眼镜、高端电视的Micro-LED显示制造",
    "",
    "3D打印电子：",
    "• 喷墨打印导电墨水、介电材料和半导体墨水",
    "• 制造柔性电路、传感器和天线",
    "• 皮升级体积控制确保打印分辨率和功能性能"
])

# 第12页：应用场景5 - 药物筛选与高通量实验
add_content_slide(prs, "应用场景五：药物筛选与高通量实验", [
    "高通量筛选（HTS）：",
    "• 每天可处理数十万个化合物样品",
    "• 皮升级分配实现纳摩尔级化合物浓度的精确配制",
    "• 大幅降低昂贵试剂和稀有样品的消耗",
    "",
    "类器官药物反应测试：",
    "• 在类器官模型上高通量测试药物反应",
    "• 皮升级分配实现多种药物组合的矩阵测试",
    "• 加速新药发现和个性化用药指导",
    "",
    "技术优势：",
    "• 提高筛选通量和数据质量",
    "• 减少实验动物使用，符合3R原则",
    "• 支持微量稀有样品的充分利用"
])

# 第13页：主要厂商对比
add_two_column_slide(prs, "主要厂商与竞争格局",
    "国际领先厂商", [
        "Eppendorf（德国）：epMotion系列自动化液体处理",
        "Tecan（瑞士）：Fluent和Freedom EVO平台",
        "Hamilton（美国）：MicroLab STAR模块化系统",
        "Beckman Coulter：Echo声波移液技术",
        "Scienion（德国）：sciFLEXARRAYER喷墨点样"
    ],
    "国产厂商崛起", [
        "华大智造：MGISP系列自动化样本制备",
        "奥美泰克：液体处理工作站国产替代",
        "迪谱诊断：数字PCR系统配套模块",
        "技术突破：在基因测序领域实现国产替代",
        "市场机遇：政策支持+供应链安全需求"
    ]
)

# 第14页：市场数据
add_content_slide(prs, "市场规模与增长预测", [
    "整体市场规模：",
    "• 2028年预计突破120亿美元",
    "• 年复合增长率（CAGR）18%",
    "• 基因测序细分市场2027年达45亿美元",
    "",
    "增长驱动因素：",
    "• 精准医疗和个性化用药需求增长",
    "• 合成生物学和基因编辑技术快速发展",
    "• 新冠疫情期间POCT市场爆发式增长",
    "• 先进封装和Micro-LED显示技术推动",
    "",
    "区域市场：",
    "• 北美和欧洲占据主要市场份额",
    "• 亚太地区增速最快，中国市场潜力巨大",
    "• 国产替代政策加速本土厂商发展"
])

# 第15页：技术挑战与未来趋势
add_content_slide(prs, "技术挑战与未来趋势", [
    "当前技术挑战：",
    "• 超高粘度液体分配：生物墨水、细胞悬浮液仍需突破",
    "• 多液体兼容性：同一系统处理有机溶剂到水溶液",
    "• 无菌与GMP合规：生物医药应用的高标准要求",
    "",
    "未来发展趋势：",
    "• 智能化与AI集成：机器学习实时优化，预测性维护",
    "• 超大规模并行化：单芯片10,000+通道，单秒百万次分配",
    "• 生物相容性突破：全生物降解系统，体内药物递送",
    "• 成本平民化：MEMS规模化制造，成本降至万元级别",
    "",
    "战略意义：",
    "• 支撑精准医疗、合成生物学和先进制造持续革新"
])

# 第16页：结语
add_title_slide(prs,
    "谢谢观看",
    "皮升级液体分配技术\n微观世界的精准操控革命")

# 保存文件
output_path = "/Users/brightfu/tech/皮升级液体分配技术.pptx"
prs.save(output_path)
print(f"PPT已生成：{output_path}")
